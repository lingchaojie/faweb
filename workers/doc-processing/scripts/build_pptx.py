import json
import sys
from io import BytesIO
from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


def points(value):
    return Inches(float(value) / 72.0)


def bbox_to_position(bbox):
    x1, y1, x2, y2 = bbox
    return points(x1), points(y1), points(max(1, x2 - x1)), points(max(1, y2 - y1))


def parse_color(value, fallback="000000"):
    if not value or not isinstance(value, str):
        value = fallback
    value = value.lstrip("#")
    if len(value) != 6:
        value = fallback
    try:
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except ValueError:
        return RGBColor(int(fallback[0:2], 16), int(fallback[2:4], 16), int(fallback[4:6], 16))


def contains_cjk(text):
    return any("㐀" <= char <= "鿿" or "豈" <= char <= "﫿" for char in text or "")


def choose_font_family(font_family, text):
    if contains_cjk(text):
        return "Noto Sans CJK SC"
    if font_family and "CID" not in font_family:
        return font_family
    return "Arial"


def set_run_font(run, font_family, text):
    run.font.name = font_family
    if contains_cjk(text):
        r_pr = run._r.get_or_add_rPr()
        ea = r_pr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(r_pr, qn("a:ea"))
        ea.set("typeface", font_family)


def visible_spans(blocks):
    spans = []
    for block in blocks:
        for span in block.get("spans", []):
            if str(span.get("text", "")).strip():
                spans.append(span)
    return spans


def style_from_source_blocks(blocks, role="body"):
    spans = visible_spans(blocks)
    largest = max(spans, key=lambda span: float(span.get("size") or 0), default={})
    text = " ".join(block.get("text", "") for block in blocks)
    fallback_size = 20 if role in {"title", "heading"} else 14
    return {
        "fontSize": float(largest.get("size") or fallback_size),
        "fontFamily": choose_font_family(largest.get("font"), text),
        "color": largest.get("color") or "#111111",
        "align": "left",
        "bullet": False,
    }


def merge_styles(source_style, hint_style):
    merged = dict(source_style)
    for key, value in (hint_style or {}).items():
        if value is not None:
            merged[key] = value
    return merged


def paragraph_alignment(value):
    if value == "center":
        return PP_ALIGN.CENTER
    if value == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def find_page_hints(hints, page_number):
    for page in hints.get("pages", []):
        if int(page.get("pageNumber", 0)) == int(page_number):
            return page
    return {"mergedTextBlocks": [], "tables": [], "regions": [], "fallbacks": [], "ignoredBlockIds": [], "imageRoles": []}


def add_text_box(slide, text_item, source_blocks=None):
    source_blocks = source_blocks or []
    text = text_item.get("text", "")
    role = text_item.get("role", "body")
    style = merge_styles(style_from_source_blocks(source_blocks, role), text_item.get("style", {}))
    left, top, width, height = bbox_to_position(text_item["bbox"])
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = paragraph_alignment(style.get("align"))
    if style.get("bullet") and not text.lstrip().startswith(("•", "-")):
        text = f"• {text}"
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(float(style.get("fontSize") or 14))
    set_run_font(run, choose_font_family(style.get("fontFamily"), text), text)
    run.font.color.rgb = parse_color(style.get("color"), "111111")
    return box


def add_drawing(slide, drawing):
    left, top, width, height = bbox_to_position(drawing["bbox"])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if drawing.get("fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = parse_color(drawing.get("fill"), "ffffff")
    else:
        shape.fill.background()
    if drawing.get("stroke"):
        shape.line.color.rgb = parse_color(drawing.get("stroke"), "000000")
    return shape


def add_image(slide, manifest_root, image):
    path = manifest_root / image["path"]
    if not path.exists():
        return None
    left, top, width, height = bbox_to_position(image["bbox"])
    return slide.shapes.add_picture(str(path), left, top, width, height)


def merged_text_from_source_blocks(item, text_blocks_by_id):
    source_blocks = [text_blocks_by_id[source_id] for source_id in item.get("sourceTextBlockIds", []) if source_id in text_blocks_by_id]
    source_blocks.sort(key=lambda block: (block.get("bbox", [0, 0, 0, 0])[1], block.get("bbox", [0, 0, 0, 0])[0]))
    return " ".join(block.get("text", "") for block in source_blocks if block.get("text", ""))


def bbox_center(bbox):
    x1, y1, x2, y2 = bbox
    return (x1 + x2) / 2, (y1 + y2) / 2


def point_in_bbox(point, bbox):
    x, y = point
    x1, y1, x2, y2 = bbox
    return x1 <= x <= x2 and y1 <= y <= y2


def bbox_overlap_ratio(bbox, candidate):
    x1, y1, x2, y2 = bbox
    cx1, cy1, cx2, cy2 = candidate
    overlap_width = max(0, min(x2, cx2) - max(x1, cx1))
    overlap_height = max(0, min(y2, cy2) - max(y1, cy1))
    bbox_area = max(0, x2 - x1) * max(0, y2 - y1)
    if bbox_area <= 0:
        return 0
    return (overlap_width * overlap_height) / bbox_area


def bbox_inside_any(bbox, bboxes):
    return any(point_in_bbox(bbox_center(bbox), candidate) or bbox_overlap_ratio(bbox, candidate) >= 0.25 for candidate in bboxes)


def fallback_candidates(page_hints):
    candidates = []
    for fallback in page_hints.get("fallbacks", []):
        candidates.append({"id": fallback.get("id", "fallback"), "bbox": fallback.get("bbox"), "zIndex": fallback.get("zIndex", 0)})
    for region in page_hints.get("regions", []):
        if region.get("strategy") == "image":
            candidates.append({"id": region.get("id", "region"), "bbox": region.get("bbox"), "zIndex": region.get("zIndex", 0)})
    return sorted([candidate for candidate in candidates if candidate.get("bbox")], key=lambda item: item.get("zIndex", 0))


def ignored_source_ids_from_regions(page_hints):
    ignored = set()
    for region in page_hints.get("regions", []):
        if region.get("strategy") in {"image", "ignore"}:
            ignored.update(region.get("sourceIds", []))
    return ignored


def crop_page_image(manifest_root, page, bbox, crop_id):
    image_path = manifest_root / page.get("imagePath", "")
    if not image_path.exists():
        return None
    image = Image.open(image_path).convert("RGB")
    scale_x = image.width / float(page["width"])
    scale_y = image.height / float(page["height"])
    x1, y1, x2, y2 = bbox
    crop_box = (
        max(0, round(x1 * scale_x)),
        max(0, round(y1 * scale_y)),
        min(image.width, round(x2 * scale_x)),
        min(image.height, round(y2 * scale_y)),
    )
    if crop_box[2] <= crop_box[0] or crop_box[3] <= crop_box[1]:
        return None
    crop_stream = BytesIO()
    safe_id = "".join(char if char.isalnum() or char in "-_" else "-" for char in str(crop_id))
    crop_stream.name = f"page-{page['pageNumber']:03d}-{safe_id}.png"
    image.crop(crop_box).save(crop_stream, format="PNG")
    crop_stream.seek(0)
    return crop_stream


def add_fallback_crop(slide, manifest_root, page, fallback):
    crop_stream = crop_page_image(manifest_root, page, fallback["bbox"], fallback["id"])
    if not crop_stream:
        return None
    left, top, width, height = bbox_to_position(fallback["bbox"])
    return slide.shapes.add_picture(crop_stream, left, top, width, height)


def table_is_editable(table_hint):
    if float(table_hint.get("confidence", 1)) < 0.75:
        return False
    rows = int(table_hint.get("rows", 0))
    columns = int(table_hint.get("columns", 0))
    return rows > 0 and columns > 0


def table_source_block_ids(table_hint, text_blocks_by_id):
    source_ids = list(table_hint.get("sourceTextBlockIds", []))
    seen = set(source_ids)
    for source_id, block in text_blocks_by_id.items():
        if source_id not in seen and bbox_inside_any(block.get("bbox", [0, 0, 0, 0]), [table_hint["bbox"]]):
            source_ids.append(source_id)
            seen.add(source_id)
    return source_ids


def add_table(slide, table_hint, text_blocks_by_id):
    if not table_is_editable(table_hint):
        return False
    rows = int(table_hint.get("rows", 0))
    columns = int(table_hint.get("columns", 0))
    left, top, width, height = bbox_to_position(table_hint["bbox"])
    table_shape = slide.shapes.add_table(rows, columns, left, top, width, height)
    table = table_shape.table
    x1, y1, x2, y2 = table_hint["bbox"]
    table_width = max(1, x2 - x1)
    table_height = max(1, y2 - y1)
    cells = {(row, column): [] for row in range(rows) for column in range(columns)}
    for source_id in table_source_block_ids(table_hint, text_blocks_by_id):
        block = text_blocks_by_id.get(source_id)
        if not block:
            continue
        cx, cy = bbox_center(block.get("bbox", table_hint["bbox"]))
        row = min(rows - 1, max(0, int(((cy - y1) / table_height) * rows)))
        column = min(columns - 1, max(0, int(((cx - x1) / table_width) * columns)))
        cells[(row, column)].append(block.get("text", ""))
    for (row, column), values in cells.items():
        table.cell(row, column).text = "\n".join(value for value in values if value)
    return True


def build_pptx(manifest_path, hints_path, output_path):
    manifest_path = Path(manifest_path)
    hints_path = Path(hints_path)
    output_path = Path(output_path)
    manifest = json.loads(manifest_path.read_text())
    hints = json.loads(hints_path.read_text())

    prs = Presentation()
    first_page = manifest["pages"][0]
    prs.slide_width = points(first_page["width"])
    prs.slide_height = points(first_page["height"])
    blank_layout = prs.slide_layouts[6]

    for page in manifest["pages"]:
        slide = prs.slides.add_slide(blank_layout)
        page_hints = find_page_hints(hints, page["pageNumber"])
        ignored = set(page_hints.get("ignoredBlockIds", [])) | ignored_source_ids_from_regions(page_hints)
        merged_source_ids = set()
        table_source_ids = set()
        fallback_boxes = [candidate["bbox"] for candidate in fallback_candidates(page_hints)]
        table_boxes = [table_hint["bbox"] for table_hint in page_hints.get("tables", []) if table_hint.get("bbox") and not bbox_inside_any(table_hint.get("bbox"), fallback_boxes) and table_is_editable(table_hint)]
        suppression_boxes = fallback_boxes + table_boxes
        text_blocks_by_id = {block.get("id"): block for block in page.get("textBlocks", []) if block.get("id")}

        fallbacks = fallback_candidates(page_hints)

        for image in page.get("images", []):
            if image.get("id") not in ignored and not bbox_inside_any(image.get("bbox", [0, 0, 0, 0]), suppression_boxes):
                add_image(slide, manifest_path.parent, image)

        for drawing in page.get("drawings", []):
            if drawing.get("id") not in ignored and not bbox_inside_any(drawing.get("bbox", [0, 0, 0, 0]), suppression_boxes):
                add_drawing(slide, drawing)

        for fallback in fallbacks:
            add_fallback_crop(slide, manifest_path.parent, page, fallback)

        for table_hint in page_hints.get("tables", []):
            if not bbox_inside_any(table_hint.get("bbox", [0, 0, 0, 0]), fallback_boxes) and add_table(slide, table_hint, text_blocks_by_id):
                table_source_ids.update(table_source_block_ids(table_hint, text_blocks_by_id))

        for item in page_hints.get("mergedTextBlocks", []):
            if bbox_inside_any(item.get("bbox", [0, 0, 0, 0]), suppression_boxes):
                merged_source_ids.update(item.get("sourceTextBlockIds", []))
                continue
            source_blocks = [text_blocks_by_id[source_id] for source_id in item.get("sourceTextBlockIds", []) if source_id in text_blocks_by_id]
            source_text = merged_text_from_source_blocks(item, text_blocks_by_id)
            if source_text:
                add_text_box(slide, {**item, "text": source_text}, source_blocks)
            merged_source_ids.update(item.get("sourceTextBlockIds", []))

        for block in page.get("textBlocks", []):
            if block.get("id") in ignored or block.get("id") in merged_source_ids or block.get("id") in table_source_ids:
                continue
            if bbox_inside_any(block.get("bbox", [0, 0, 0, 0]), suppression_boxes):
                continue
            add_text_box(slide, {
                "text": block.get("text", ""),
                "bbox": block.get("bbox", [0, 0, 1, 1]),
                "role": "body",
            }, [block])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main():
    if len(sys.argv) != 4:
        raise SystemExit("Usage: build_pptx.py <manifest.json> <hints.json> <output.pptx>")
    build_pptx(sys.argv[1], sys.argv[2], sys.argv[3])


if __name__ == "__main__":
    main()
