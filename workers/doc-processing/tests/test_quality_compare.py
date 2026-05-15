import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

import fitz
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from scripts.quality_compare import (
    collect_pptx_stats,
    compare_images,
    find_sample_pairs,
    render_pdf_pages,
    render_pptx_pages,
    run_comparison,
    write_side_by_side,
)


class QualityCompareTest(unittest.TestCase):
    def test_find_sample_pairs_only_returns_pdfs_with_matching_pptx(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            (root / "alpha.pdf").write_bytes(b"%PDF-1.4\n")
            (root / "alpha.pptx").write_bytes(b"pptx")
            (root / "beta.pdf").write_bytes(b"%PDF-1.4\n")
            (root / "notes.txt").write_text("ignored")

            pairs = find_sample_pairs(root)

            self.assertEqual(len(pairs), 1)
            self.assertEqual(pairs[0].name, "alpha")
            self.assertEqual(pairs[0].pdf_path, root / "alpha.pdf")
            self.assertEqual(pairs[0].reference_pptx_path, root / "alpha.pptx")

    def test_collect_pptx_stats_counts_editable_objects(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image_path = root / "image.png"
            Image.new("RGB", (20, 20), "red").save(image_path)

            pptx_path = root / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
            textbox.text = "Editable title"
            slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.2), Inches(1), Inches(1))
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.2), Inches(1), Inches(1))
            table = slide.shapes.add_table(2, 2, Inches(0.5), Inches(2.5), Inches(2), Inches(1)).table
            table.cell(0, 0).text = "A"
            table.cell(0, 1).text = "B"
            table.cell(1, 0).text = "C"
            table.cell(1, 1).text = "D"
            prs.save(pptx_path)

            stats = collect_pptx_stats(pptx_path)

            self.assertEqual(stats["slideCount"], 1)
            self.assertEqual(stats["textShapes"], 1)
            self.assertEqual(stats["pictures"], 1)
            self.assertGreaterEqual(stats["shapes"], 2)
            self.assertEqual(stats["tables"], 1)

    def test_render_pdf_pages_writes_numbered_pngs(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            pdf_path = root / "input.pdf"
            doc = fitz.open()
            page = doc.new_page(width=200, height=100)
            page.insert_text((20, 40), "Hello")
            doc.save(pdf_path)
            doc.close()

            pages = render_pdf_pages(pdf_path, root / "renders", zoom=1)

            self.assertEqual([path.name for path in pages], ["page-001.png"])
            self.assertTrue(pages[0].exists())

    def test_compare_images_returns_zero_for_identical_images(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            first = root / "first.png"
            second = root / "second.png"
            Image.new("RGB", (10, 10), "white").save(first)
            Image.new("RGB", (10, 10), "white").save(second)

            self.assertEqual(compare_images(first, second), 0.0)

    def test_compare_images_returns_positive_value_for_different_images(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            first = root / "first.png"
            second = root / "second.png"
            Image.new("RGB", (10, 10), "white").save(first)
            Image.new("RGB", (10, 10), "black").save(second)

            self.assertGreater(compare_images(first, second), 0.9)

    def test_write_side_by_side_combines_images_with_labels(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            left = root / "left.png"
            right = root / "right.png"
            output = root / "review.png"
            Image.new("RGB", (20, 10), "red").save(left)
            Image.new("RGB", (20, 10), "blue").save(right)

            write_side_by_side([("Left", left), ("Right", right)], output)

            with Image.open(output) as result:
                self.assertEqual(result.size, (40, 34))

    def test_render_pptx_pages_converts_to_pdf_then_renders(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            pptx_path = root / "deck.pptx"
            pptx_path.write_bytes(b"fake pptx")

            def fake_run(args, check, stdout, stderr):
                out_dir = Path(args[args.index("--outdir") + 1])
                doc = fitz.open()
                doc.new_page(width=200, height=100)
                doc.save(out_dir / "deck.pdf")
                doc.close()
                return subprocess.CompletedProcess(args, 0)

            with patch("scripts.quality_compare.subprocess.run", side_effect=fake_run) as run:
                pages = render_pptx_pages(pptx_path, root / "pptx-renders", libreoffice_bin="soffice", zoom=1)

            self.assertEqual(len(pages), 1)
            self.assertTrue(pages[0].exists())
            run.assert_called_once()

    def write_minimal_pptx(self, path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.5)).text = "A"
        prs.save(path)

    def test_run_comparison_writes_summary_and_review_images(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            samples = root / "samples"
            generated = root / "generated"
            output = root / "analysis"
            samples.mkdir()
            generated.mkdir()
            (samples / "alpha.pdf").write_bytes(b"%PDF-1.4\n")
            self.write_minimal_pptx(samples / "alpha.pptx")
            self.write_minimal_pptx(generated / "alpha.pptx")

            def fake_render_pdf(_source, out_dir, zoom=2.0):
                out_dir.mkdir(parents=True, exist_ok=True)
                image = out_dir / "page-001.png"
                Image.new("RGB", (20, 10), "white").save(image)
                return [image]

            def fake_render_pptx(source, out_dir, libreoffice_bin="soffice", zoom=2.0):
                out_dir.mkdir(parents=True, exist_ok=True)
                image = out_dir / "page-001.png"
                color = "white" if source.parent.name == "samples" else "black"
                Image.new("RGB", (20, 10), color).save(image)
                return [image]

            report = run_comparison(
                samples_dir=samples,
                output_dir=output,
                generated_dir=generated,
                render_pdf=fake_render_pdf,
                render_pptx=fake_render_pptx,
            )

            summary = json.loads((output / "summary.json").read_text())
            self.assertEqual(summary["samples"][0]["name"], "alpha")
            self.assertEqual(report["samples"][0]["slides"][0]["pageNumber"], 1)
            self.assertGreater(report["samples"][0]["slides"][0]["generatedDiffToReference"], 0.9)
            self.assertTrue((output / "alpha" / "review" / "page-001.png").exists())


if __name__ == "__main__":
    unittest.main()
